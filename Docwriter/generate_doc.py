import os
import fnmatch
import argparse
import re

# --- Helper Function to Sanitize Text ---
def sanitize_text(text):
    """
    Remove control characters (except standard whitespace) that are not XML‐compatible.
    """
    # Remove control characters in the ranges 0x00-0x08, 0x0B-0x0C, 0x0E-0x1F and 0x7F.
    return re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)

# --- Helper Function to Check if a File is Text ---
def is_text_file(file_path, blocksize=512):
    """
    Returns True if the file appears to be a text file based on its extension or content.
    Explicitly treat certain file types (like .pptx, .docx, .pdf, etc.) as needing special parsing.
    """
    text_extensions = ['.txt', '.py', '.md', '.json', '.csv', '.xml', '.html', '.css', '.js']
    ext = os.path.splitext(file_path)[1].lower()

    if ext in text_extensions:
        return True
    elif ext in ['.pptx', '.docx', '.xlsx', '.pdf']:
        # These file types need special parsing; do not use generic text reading
        return False
    else:
        try:
            with open(file_path, 'rb') as f:
                chunk = f.read(blocksize)
            # If a null byte is found, likely a binary file.
            return b'\0' not in chunk
        except Exception:
            return False

# --- Helper Function to Extract Text from PPTX Files ---
from pptx import Presentation
def extract_text_from_pptx(file_path):
    """
    Extract text content from a PowerPoint (.pptx) file.
    """
    text_runs = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
    return "\n".join(text_runs)

# --- Common Functions ---
def load_ignore_patterns(root_path):
    """
    Load ignore patterns from a .docignore file in the root folder, if it exists.
    Each non-blank, non-comment line is treated as an ignore pattern.
    """
    ignore_file = os.path.join(root_path, '.docignore')
    patterns = []
    if os.path.exists(ignore_file):
        with open(ignore_file, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if line and not line.startswith('#'):
                    patterns.append(line)
    return patterns

def should_ignore(relative_path, ignore_patterns):
    """
    Return True if the relative_path matches any of the ignore patterns.
    Normalizes path separators and supports wildcard patterns.
    """
    relative_path = relative_path.replace('\\', '/')
    path_parts = relative_path.split('/')
    for pattern in ignore_patterns:
        if pattern.endswith('/'):
            directory = pattern.rstrip('/')
            if directory in path_parts:
                return True
        if fnmatch.fnmatch(relative_path, pattern) or fnmatch.fnmatch(os.path.basename(relative_path), pattern):
            return True
    return False

def get_structure_lines(root_path, ignore_patterns):
    """
    Recursively builds a list of lines representing the folder structure.
    Directories end with a slash, and files are listed as plain names.
    The structure is produced in a flat order.
    """
    lines = []
    root_name = os.path.basename(os.path.abspath(root_path))
    lines.append(f"{root_name}/")
    
    items = sorted(os.listdir(root_path))
    
    # Process files first.
    for item in items:
        full_item = os.path.join(root_path, item)
        rel_item = os.path.relpath(full_item, root_path).replace('\\', '/')
        if should_ignore(rel_item, ignore_patterns):
            continue
        if os.path.isfile(full_item):
            lines.append(item)
    
    # Process directories.
    for item in items:
        full_item = os.path.join(root_path, item)
        rel_item = os.path.relpath(full_item, root_path).replace('\\', '/')
        if should_ignore(rel_item, ignore_patterns):
            continue
        if os.path.isdir(full_item):
            lines.append(f"{item}/")
            sub_lines = get_structure_lines(full_item, ignore_patterns)
            if sub_lines:
                sub_lines = sub_lines[1:]  # Remove duplicate directory name.
            lines.extend(sub_lines)
    return lines

# --- Word (.docx) Generation ---
from docx import Document
from docx.shared import Pt

def create_word_from_folder(root_path, output_file):
    ignore_patterns = load_ignore_patterns(root_path)
    print("Ignore patterns:", ignore_patterns)
    
    document = Document()
    
    # SECTION 1: Folder Structure
    document.add_heading('Folder Structure', level=1)
    structure_lines = get_structure_lines(root_path, ignore_patterns)
    for line in structure_lines:
        document.add_paragraph(line)
    
    # SECTION 2: File Paths and Contents
    document.add_page_break()
    document.add_heading('File Paths and Contents', level=1)
    
    for dirpath, dirnames, filenames in os.walk(root_path):
        dirnames[:] = [d for d in dirnames if not should_ignore(os.path.relpath(os.path.join(dirpath, d), root_path).replace('\\', '/'), ignore_patterns)]
        for filename in sorted(filenames):
            file_path = os.path.join(dirpath, filename)
            rel_file_path = os.path.relpath(file_path, root_path).replace('\\', '/')
            if should_ignore(rel_file_path, ignore_patterns):
                continue
            print(f"Processing: {file_path}")
            document.add_heading(file_path, level=2)
            try:
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.pptx':
                    file_contents = extract_text_from_pptx(file_path)
                elif not is_text_file(file_path):
                    file_contents = "[Binary file - content not displayed]"
                else:
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        file_contents = f.read()
                # Sanitize extracted text
                file_contents = sanitize_text(file_contents)
            except Exception as e:
                file_contents = f"Error reading file: {e}"
            paragraph = document.add_paragraph(file_contents)
            for run in paragraph.runs:
                run.font.name = 'Courier New'
                run.font.size = Pt(10)
            print(f"Done processing: {file_path}")
    
    document.save(output_file)
    print(f"\nWord document '{output_file}' created successfully.")

# --- PDF Generation ---
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

def create_pdf_from_folder(root_path, output_file):
    ignore_patterns = load_ignore_patterns(root_path)
    print("Ignore patterns:", ignore_patterns)
    
    c = canvas.Canvas(output_file, pagesize=letter)
    width, height = letter
    margin = 40
    y = height - margin
    line_height = 12

    def draw_line(text, y_position):
        c.drawString(margin, y_position, text)

    # SECTION 1: Folder Structure
    c.setFont("Helvetica-Bold", 14)
    draw_line("Folder Structure", y)
    y -= line_height * 2

    c.setFont("Helvetica", 10)
    structure_lines = get_structure_lines(root_path, ignore_patterns)
    for line in structure_lines:
        if y < margin:
            c.showPage()
            y = height - margin
            c.setFont("Helvetica", 10)
        draw_line(line, y)
        y -= line_height

    # SECTION 2: File Paths and Contents
    if y < margin + 40:
        c.showPage()
        y = height - margin
    else:
        y -= line_height * 2
    c.setFont("Helvetica-Bold", 14)
    draw_line("File Paths and Contents", y)
    y -= line_height * 2
    c.setFont("Courier", 8)

    for dirpath, dirnames, filenames in os.walk(root_path):
        dirnames[:] = [d for d in dirnames if not should_ignore(os.path.relpath(os.path.join(dirpath, d), root_path).replace('\\', '/'), ignore_patterns)]
        for filename in sorted(filenames):
            file_path = os.path.join(dirpath, filename)
            rel_file_path = os.path.relpath(file_path, root_path).replace('\\', '/')
            if should_ignore(rel_file_path, ignore_patterns):
                continue
            heading = f"File: {file_path}"
            if y < margin + line_height * 2:
                c.showPage()
                y = height - margin
            c.setFont("Helvetica-Bold", 10)
            draw_line(heading, y)
            y -= line_height
            c.setFont("Courier", 8)
            try:
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.pptx':
                    file_contents = extract_text_from_pptx(file_path)
                elif not is_text_file(file_path):
                    file_contents = "[Binary file - content not displayed]"
                else:
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        file_contents = f.read()
                file_contents = sanitize_text(file_contents)
            except Exception as e:
                file_contents = f"Error reading file: {e}"
            for content_line in file_contents.splitlines():
                if y < margin:
                    c.showPage()
                    y = height - margin
                    c.setFont("Courier", 8)
                draw_line(content_line, y)
                y -= line_height
            y -= line_height  # extra space between files
            print(f"Done processing: {file_path}")

    c.save()
    print(f"\nPDF file '{output_file}' created successfully.")

# --- PowerPoint (.pptx) Generation ---
from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt_from_folder(root_path, output_file):
    ignore_patterns = load_ignore_patterns(root_path)
    print("Ignore patterns:", ignore_patterns)
    
    prs = Presentation()

    # Slide 1: Folder Structure
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = "Folder Structure"
    
    structure_lines = get_structure_lines(root_path, ignore_patterns)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    content_tf = content_box.text_frame
    for line in structure_lines:
        p = content_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)

    # Slide 2: File Paths and Contents overview
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text = "File Paths and Contents"

    for dirpath, dirnames, filenames in os.walk(root_path):
        dirnames[:] = [d for d in dirnames if not should_ignore(os.path.relpath(os.path.join(dirpath, d), root_path).replace('\\', '/'), ignore_patterns)]
        for filename in sorted(filenames):
            file_path = os.path.join(dirpath, filename)
            rel_file_path = os.path.relpath(file_path, root_path).replace('\\', '/')
            if should_ignore(rel_file_path, ignore_patterns):
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_tf = title_box.text_frame
            title_tf.text = file_path
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
            content_tf = content_box.text_frame
            try:
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.pptx':
                    file_contents = extract_text_from_pptx(file_path)
                elif not is_text_file(file_path):
                    file_contents = "[Binary file - content not displayed]"
                else:
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        file_contents = f.read()
                file_contents = sanitize_text(file_contents)
            except Exception as e:
                file_contents = f"Error reading file: {e}"
            for line in file_contents.splitlines():
                p = content_tf.add_paragraph()
                p.text = line
                p.font.size = Pt(10)
            print(f"Done processing: {file_path}")
            
    prs.save(output_file)
    print(f"\nPowerPoint file '{output_file}' created successfully.")

# --- Jupyter Notebook (.ipynb) Generation ---
import nbformat

def create_ipynb_from_folder(root_path, output_file):
    ignore_patterns = load_ignore_patterns(root_path)
    print("Ignore patterns:", ignore_patterns)
    
    nb = nbformat.v4.new_notebook()
    cells = []

    # Cell 1: Folder Structure
    structure_lines = get_structure_lines(root_path, ignore_patterns)
    md_structure = "# Folder Structure\n\n```\n" + "\n".join(structure_lines) + "\n```"
    cells.append(nbformat.v4.new_markdown_cell(md_structure))
    
    # Cell 2: File Paths and Contents overview
    cells.append(nbformat.v4.new_markdown_cell("# File Paths and Contents"))

    for dirpath, dirnames, filenames in os.walk(root_path):
        dirnames[:] = [d for d in dirnames if not should_ignore(os.path.relpath(os.path.join(dirpath, d), root_path).replace('\\', '/'), ignore_patterns)]
        for filename in sorted(filenames):
            file_path = os.path.join(dirpath, filename)
            rel_file_path = os.path.relpath(file_path, root_path).replace('\\', '/')
            if should_ignore(rel_file_path, ignore_patterns):
                continue
            try:
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.pptx':
                    file_contents = extract_text_from_pptx(file_path)
                elif not is_text_file(file_path):
                    file_contents = "[Binary file - content not displayed]"
                else:
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        file_contents = f.read()
                file_contents = sanitize_text(file_contents)
            except Exception as e:
                file_contents = f"Error reading file: {e}"
            md_cell = f"## {file_path}\n\n```\n{file_contents}\n```"
            cells.append(nbformat.v4.new_markdown_cell(md_cell))
            print(f"Done processing: {file_path}")
    
    nb['cells'] = cells
    with open(output_file, 'w', encoding='utf-8') as f:
        nbformat.write(nb, f)
    print(f"\nJupyter Notebook '{output_file}' created successfully.")

# --- Main Function ---
def main():
    parser = argparse.ArgumentParser(
        description="Generate documentation (Word, PDF, PowerPoint, or Jupyter Notebook) from folder structure and file contents, supporting any file type."
    )
    parser.add_argument("folder", help="Folder path to scan")
    parser.add_argument("output", help="Output file with extension (.docx, .pdf, .pptx, .ipynb)")
    args = parser.parse_args()
    
    if not os.path.exists(args.folder):
        parser.error(f"Folder not found: {args.folder}")
    
    ext = os.path.splitext(args.output)[1].lower()
    if ext == ".docx":
        create_word_from_folder(args.folder, args.output)
    elif ext == ".pdf":
        create_pdf_from_folder(args.folder, args.output)
    elif ext == ".pptx":
        create_ppt_from_folder(args.folder, args.output)
    elif ext == ".ipynb":
        create_ipynb_from_folder(args.folder, args.output)
    else:
        print("Unsupported output file type. Please use one of: .docx, .pdf, .pptx, or .ipynb")

if __name__ == '__main__':
    main()
